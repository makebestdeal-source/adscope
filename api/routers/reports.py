"""One-click advertiser report bundle exports."""

import io
import re
import zipfile
from datetime import datetime, timedelta, timezone
from pathlib import Path
from urllib.parse import quote

from fastapi import APIRouter, Depends, HTTPException, Query
from fastapi.responses import StreamingResponse
from openpyxl import Workbook
from openpyxl.drawing.image import Image as XLImage
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import Image as RLImage
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
from sqlalchemy import func, select
from sqlalchemy.ext.asyncio import AsyncSession

from api.deps import get_current_user, require_paid
from database import get_db
from database.models import (
    AdDetail,
    AdSnapshot,
    Advertiser,
    BrandChannelContent,
    Campaign,
    Industry,
    SpendEstimate,
    User,
)

router = APIRouter(
    prefix="/api/reports",
    tags=["reports"],
    dependencies=[Depends(require_paid)],
)

KST = timezone(timedelta(hours=9))
BRAND = RGBColor(20, 184, 166)
INK = RGBColor(15, 23, 42)
MUTED = RGBColor(100, 116, 139)

_HEADER_FILL = PatternFill(start_color="0F172A", end_color="0F172A", fill_type="solid")
_HEADER_FONT = Font(bold=True, color="FFFFFF", size=11)
_HEADER_ALIGN = Alignment(horizontal="center", vertical="center")

_CHANNEL_LABELS = {
    "naver_search": "네이버 검색",
    "naver_da": "네이버 DA",
    "naver_shopping": "네이버 쇼핑",
    "kakao_da": "카카오 DA",
    "youtube_ads": "YouTube",
    "meta": "Meta",
    "tiktok_ads": "TikTok",
    "google_gdn": "Google GDN",
    "google_search_ads": "Google 검색",
}


def _now_kst() -> datetime:
    return datetime.now(KST)


def _default_range(days: int) -> tuple[datetime, datetime]:
    date_to = datetime.utcnow()
    return date_to - timedelta(days=days), date_to


def _safe(value) -> str:
    return "" if value is None else str(value)


def _money(value: float | int | None) -> str:
    value = float(value or 0)
    if value >= 100_000_000:
        return f"{value / 100_000_000:.1f}억"
    if value >= 10_000:
        return f"{value / 10_000:.0f}만"
    return f"{value:,.0f}"


def _date(value: datetime | None) -> str:
    if value is None:
        return ""
    return (value + timedelta(hours=9)).strftime("%Y-%m-%d")


def _channel(value: str | None) -> str:
    return _CHANNEL_LABELS.get(value or "", value or "-")


def _filename_safe(value: str) -> str:
    cleaned = re.sub(r'[\\/:*?"<>|]+', "_", value or "advertiser")
    return cleaned.strip(" ._")[:80] or "advertiser"


def _content_disposition(filename: str) -> dict[str, str]:
    ascii_name = filename.encode("ascii", "replace").decode("ascii").replace("?", "_")
    utf8_name = quote(filename, safe="")
    return {"Content-Disposition": f"attachment; filename=\"{ascii_name}\"; filename*=UTF-8''{utf8_name}"}


def _resolve_image_path(raw_path: str | None) -> str | None:
    if not raw_path:
        return None

    clean = raw_path.replace("\\", "/").split("#", 1)[0].lstrip("/")
    if ".." in clean:
        return None

    root = Path.cwd().resolve()
    image_root = Path("stored_images").resolve()
    screenshot_root = Path("screenshots").resolve()
    candidates: list[Path] = []

    candidates.append((root / clean).resolve())
    if clean.startswith("stored_images/"):
        candidates.append((image_root / clean[len("stored_images/"):]).resolve())
    else:
        candidates.append((image_root / clean).resolve())

    if clean.startswith("screenshots/"):
        candidates.append((screenshot_root / clean[len("screenshots/"):]).resolve())
    else:
        candidates.append((screenshot_root / clean).resolve())

    for candidate in candidates:
        try:
            candidate.relative_to(root)
        except ValueError:
            continue
        if candidate.exists() and candidate.is_file():
            return str(candidate)
    return None


def _register_pdf_font() -> str:
    candidates = [
        Path("C:/Windows/Fonts/malgun.ttf"),
        Path("/usr/share/fonts/truetype/nanum/NanumGothic.ttf"),
        Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
    ]
    for path in candidates:
        if path.exists():
            try:
                pdfmetrics.registerFont(TTFont("AdScopeKR", str(path)))
                return "AdScopeKR"
            except Exception:
                continue
    return "Helvetica"


async def _collect_report_data(
    db: AsyncSession,
    advertiser_id: int,
    days: int,
) -> dict:
    advertiser = (
        await db.execute(select(Advertiser).where(Advertiser.id == advertiser_id))
    ).scalar_one_or_none()
    if not advertiser:
        raise HTTPException(status_code=404, detail="Advertiser not found")

    date_from, date_to = _default_range(days)

    industry_name = ""
    if advertiser.industry_id:
        industry = (
            await db.execute(select(Industry).where(Industry.id == advertiser.industry_id))
        ).scalar_one_or_none()
        industry_name = industry.name if industry else ""

    channel_rows = (
        await db.execute(
            select(
                AdSnapshot.channel,
                func.count(AdDetail.id).label("ad_count"),
                func.min(AdSnapshot.captured_at).label("first_seen"),
                func.max(AdSnapshot.captured_at).label("last_seen"),
            )
            .join(AdSnapshot, AdDetail.snapshot_id == AdSnapshot.id)
            .where(AdDetail.advertiser_id == advertiser_id)
            .where(AdSnapshot.captured_at >= date_from)
            .where(AdSnapshot.captured_at <= date_to)
            .group_by(AdSnapshot.channel)
            .order_by(func.count(AdDetail.id).desc())
        )
    ).all()

    spend_rows = (
        await db.execute(
            select(
                SpendEstimate.channel,
                func.sum(SpendEstimate.est_daily_spend).label("total_spend"),
            )
            .join(Campaign, SpendEstimate.campaign_id == Campaign.id)
            .where(Campaign.advertiser_id == advertiser_id)
            .where(SpendEstimate.date >= date_from)
            .where(SpendEstimate.date <= date_to)
            .group_by(SpendEstimate.channel)
        )
    ).all()
    spend_by_channel = {r.channel: float(r.total_spend or 0) for r in spend_rows}

    channels = [
        {
            "channel": r.channel,
            "channel_label": _channel(r.channel),
            "ad_count": int(r.ad_count or 0),
            "first_seen": r.first_seen,
            "last_seen": r.last_seen,
            "spend": spend_by_channel.get(r.channel, 0),
        }
        for r in channel_rows
    ]

    campaign_rows = (
        await db.execute(
            select(
                Campaign.campaign_name,
                Campaign.channel,
                Campaign.first_seen,
                Campaign.last_seen,
                Campaign.is_active,
                Campaign.total_est_spend,
                Campaign.snapshot_count,
                Campaign.objective,
                Campaign.product_service,
            )
            .where(Campaign.advertiser_id == advertiser_id)
            .where(Campaign.last_seen >= date_from)
            .order_by(Campaign.last_seen.desc())
            .limit(12)
        )
    ).all()

    campaigns = [
        {
            "name": r.campaign_name or "캠페인",
            "channel": _channel(r.channel),
            "first_seen": r.first_seen,
            "last_seen": r.last_seen,
            "active": bool(r.is_active),
            "spend": float(r.total_est_spend or 0),
            "snapshots": int(r.snapshot_count or 0),
            "objective": r.objective or "",
            "product_service": r.product_service or "",
        }
        for r in campaign_rows
        if not re.search(r"라이브러리\s*ID\s*:|library\s*id\s*:", r.campaign_name or "", re.I)
    ]

    creative_rows = (
        await db.execute(
            select(
                AdSnapshot.captured_at,
                AdSnapshot.channel,
                AdDetail.ad_text,
                AdDetail.ad_description,
                AdDetail.product_name,
                AdDetail.creative_image_path,
                AdDetail.screenshot_path,
                AdDetail.url,
            )
            .join(AdSnapshot, AdDetail.snapshot_id == AdSnapshot.id)
            .where(AdDetail.advertiser_id == advertiser_id)
            .where(AdSnapshot.captured_at >= date_from)
            .where(AdSnapshot.captured_at <= date_to)
            .order_by(AdSnapshot.captured_at.desc())
            .limit(30)
        )
    ).all()

    creatives = [
        {
            "date": r.captured_at,
            "channel": _channel(r.channel),
            "text": (r.ad_text or r.ad_description or "")[:180],
            "product": r.product_name or "",
            "image_path": _resolve_image_path(r.creative_image_path or r.screenshot_path),
            "url": r.url or "",
        }
        for r in creative_rows
    ]
    image_creatives = [creative for creative in creatives if creative["image_path"]]

    social_rows = (
        await db.execute(
            select(
                BrandChannelContent.platform,
                BrandChannelContent.upload_date,
                BrandChannelContent.title,
                BrandChannelContent.content_type,
                BrandChannelContent.view_count,
                BrandChannelContent.like_count,
            )
            .where(BrandChannelContent.advertiser_id == advertiser_id)
            .order_by(BrandChannelContent.upload_date.desc())
            .limit(12)
        )
    ).all()
    social = [
        {
            "platform": r.platform or "",
            "date": r.upload_date,
            "title": r.title or "",
            "type": r.content_type or "",
            "views": int(r.view_count or 0),
            "likes": int(r.like_count or 0),
        }
        for r in social_rows
    ]

    total_ads = sum(c["ad_count"] for c in channels)
    total_spend = sum(c["spend"] for c in channels) or sum(c["spend"] for c in campaigns)
    active_campaigns = sum(1 for c in campaigns if c["active"])
    top_channel = channels[0]["channel_label"] if channels else "-"

    insights = []
    if total_ads:
        insights.append(f"최근 {days}일 동안 {total_ads:,}건의 광고 소재가 관측되었습니다.")
    if top_channel != "-":
        insights.append(f"가장 활발한 채널은 {top_channel}입니다.")
    if total_spend:
        insights.append(f"동기간 추정 광고비 합계는 약 {_money(total_spend)}원입니다.")
    if active_campaigns:
        insights.append(f"진행 중으로 분류된 캠페인은 {active_campaigns}건입니다.")
    if not insights:
        insights.append("선택 기간에 보고서로 정리할 수 있는 광고 활동 데이터가 부족합니다.")

    return {
        "advertiser": advertiser,
        "industry_name": industry_name,
        "date_from": date_from,
        "date_to": date_to,
        "days": days,
        "channels": channels,
        "campaigns": campaigns,
        "creatives": creatives,
        "image_creatives": image_creatives,
        "social": social,
        "summary": {
            "total_ads": total_ads,
            "total_spend": total_spend,
            "campaign_count": len(campaigns),
            "active_campaigns": active_campaigns,
            "top_channel": top_channel,
            "insights": insights[:5],
        },
    }


def _style_header(ws, count: int):
    for col in range(1, count + 1):
        cell = ws.cell(row=1, column=col)
        cell.fill = _HEADER_FILL
        cell.font = _HEADER_FONT
        cell.alignment = _HEADER_ALIGN


def _auto_width(ws, count: int):
    for col in range(1, count + 1):
        width = 10
        for row in ws.iter_rows(min_col=col, max_col=col):
            for cell in row:
                width = max(width, min(len(str(cell.value or "")) + 2, 48))
        ws.column_dimensions[get_column_letter(col)].width = width


def _build_xlsx(data: dict) -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = "요약"
    rows = [
        ["항목", "값"],
        ["광고주", data["advertiser"].name],
        ["업종", data["industry_name"] or "-"],
        ["분석 기간", f"{_date(data['date_from'])} ~ {_date(data['date_to'])}"],
        ["광고 소재 수", data["summary"]["total_ads"]],
        ["캠페인 수", data["summary"]["campaign_count"]],
        ["진행 중 캠페인", data["summary"]["active_campaigns"]],
        ["주요 채널", data["summary"]["top_channel"]],
        ["추정 광고비", data["summary"]["total_spend"]],
    ]
    for row in rows:
        ws.append(row)
    _style_header(ws, 2)
    _auto_width(ws, 2)

    ws2 = wb.create_sheet("채널 요약")
    headers = ["채널", "소재 수", "최초 관측", "최근 관측", "추정 광고비"]
    ws2.append(headers)
    for r in data["channels"]:
        ws2.append([r["channel_label"], r["ad_count"], _date(r["first_seen"]), _date(r["last_seen"]), r["spend"]])
    _style_header(ws2, len(headers))
    _auto_width(ws2, len(headers))

    ws3 = wb.create_sheet("캠페인")
    headers = ["캠페인", "채널", "상태", "최초 관측", "최근 관측", "소재/스냅샷", "추정 광고비", "목적", "제품/서비스"]
    ws3.append(headers)
    for r in data["campaigns"]:
        ws3.append([
            r["name"], r["channel"], "진행 중" if r["active"] else "종료",
            _date(r["first_seen"]), _date(r["last_seen"]), r["snapshots"], r["spend"],
            r["objective"], r["product_service"],
        ])
    _style_header(ws3, len(headers))
    _auto_width(ws3, len(headers))

    ws4 = wb.create_sheet("소재")
    headers = ["관측일", "채널", "제품", "광고 문구", "URL"]
    ws4.append(headers)
    for r in data["creatives"]:
        ws4.append([_date(r["date"]), r["channel"], r["product"], r["text"], r["url"]])
    _style_header(ws4, len(headers))
    _auto_width(ws4, len(headers))

    ws_img = wb.create_sheet("이미지 소재")
    headers = ["이미지", "관측일", "채널", "제품", "광고 문구", "URL"]
    ws_img.append(headers)
    _style_header(ws_img, len(headers))
    ws_img.column_dimensions["A"].width = 24
    for idx, r in enumerate(data["image_creatives"][:20], start=2):
        ws_img.cell(row=idx, column=2, value=_date(r["date"]))
        ws_img.cell(row=idx, column=3, value=r["channel"])
        ws_img.cell(row=idx, column=4, value=r["product"])
        ws_img.cell(row=idx, column=5, value=r["text"])
        ws_img.cell(row=idx, column=6, value=r["url"])
        ws_img.row_dimensions[idx].height = 82
        try:
            image = XLImage(r["image_path"])
            image.width = 140
            image.height = 96
            ws_img.add_image(image, f"A{idx}")
        except Exception:
            ws_img.cell(row=idx, column=1, value="이미지 로드 실패")
    _auto_width(ws_img, len(headers))
    ws_img.column_dimensions["A"].width = 24

    ws5 = wb.create_sheet("소셜")
    headers = ["플랫폼", "게시일", "유형", "제목", "조회수", "좋아요"]
    ws5.append(headers)
    for r in data["social"]:
        ws5.append([r["platform"], _date(r["date"]), r["type"], r["title"], r["views"], r["likes"]])
    _style_header(ws5, len(headers))
    _auto_width(ws5, len(headers))

    for sheet in wb.worksheets:
        sheet.freeze_panes = "A2"
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_picture_fit(slide, path: str, x: float, y: float, w: float, h: float):
    try:
        with PILImage.open(path) as image:
            img_w, img_h = image.size
            converted = image.convert("RGB")
            image_buf = io.BytesIO()
            converted.save(image_buf, format="PNG")
            image_buf.seek(0)
        if not img_w or not img_h:
            return None
        scale = min(w / img_w, h / img_h)
        draw_w = img_w * scale
        draw_h = img_h * scale
        return slide.shapes.add_picture(
            image_buf,
            Inches(x + (w - draw_w) / 2),
            Inches(y + (h - draw_h) / 2),
            width=Inches(draw_w),
            height=Inches(draw_h),
        )
    except Exception:
        return None


def _build_pptx(data: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    _add_textbox(slide, 0.7, 0.55, 2.0, 0.35, "AdScope", 15, True, BRAND)
    _add_textbox(slide, 0.7, 1.55, 8.8, 0.8, f"{data['advertiser'].name} 광고 활동 리포트", 32, True)
    _add_textbox(slide, 0.72, 2.45, 6.0, 0.35, f"{_date(data['date_from'])} ~ {_date(data['date_to'])}", 15, False, MUTED)
    _add_textbox(slide, 0.72, 6.75, 6.0, 0.3, f"Generated {_now_kst().strftime('%Y-%m-%d %H:%M KST')}", 10, False, MUTED)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, 0.55, 0.35, 4.5, 0.4, "Executive Summary", 22, True)
    metrics = [
        ("광고 소재", f"{data['summary']['total_ads']:,}건"),
        ("캠페인", f"{data['summary']['campaign_count']:,}건"),
        ("추정 광고비", f"{_money(data['summary']['total_spend'])}원"),
        ("주요 채널", data["summary"]["top_channel"]),
    ]
    for idx, (label, value) in enumerate(metrics):
        x = 0.6 + idx * 3.1
        shape = slide.shapes.add_shape(1, Inches(x), Inches(1.15), Inches(2.75), Inches(1.25))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(241, 245, 249)
        shape.line.color.rgb = RGBColor(226, 232, 240)
        _add_textbox(slide, x + 0.18, 1.35, 2.3, 0.25, label, 11, False, MUTED)
        _add_textbox(slide, x + 0.18, 1.72, 2.3, 0.42, value, 19, True, INK)
    for i, insight in enumerate(data["summary"]["insights"]):
        _add_textbox(slide, 0.85, 3.1 + i * 0.48, 11.5, 0.32, f"- {insight}", 15)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, 0.55, 0.35, 5.0, 0.4, "채널별 집행 요약", 22, True)
    table_rows = min(len(data["channels"]), 8) + 1
    table = slide.shapes.add_table(table_rows, 5, Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.42 * table_rows)).table
    headers = ["채널", "소재 수", "최초 관측", "최근 관측", "추정 광고비"]
    for col, header in enumerate(headers):
        table.cell(0, col).text = header
    for row_idx, row in enumerate(data["channels"][:8], start=1):
        values = [row["channel_label"], f"{row['ad_count']:,}", _date(row["first_seen"]), _date(row["last_seen"]), f"{_money(row['spend'])}원"]
        for col, value in enumerate(values):
            table.cell(row_idx, col).text = value
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Malgun Gothic"
                paragraph.font.size = Pt(10)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, 0.55, 0.35, 5.0, 0.4, "주요 광고 소재", 22, True)
    image_items = data["image_creatives"][:6]
    if image_items:
        for i, creative in enumerate(image_items):
            col = i % 3
            row = i // 3
            x = 0.7 + col * 4.15
            y = 1.05 + row * 3.0
            bg = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.55), Inches(2.45))
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(248, 250, 252)
            bg.line.color.rgb = RGBColor(226, 232, 240)
            _add_picture_fit(slide, creative["image_path"], x + 0.12, y + 0.12, 3.31, 1.7)
            _add_textbox(slide, x + 0.15, y + 1.9, 3.25, 0.2, f"{_date(creative['date'])} | {creative['channel']}", 8, False, MUTED)
            _add_textbox(slide, x + 0.15, y + 2.12, 3.25, 0.25, creative["text"] or "-", 9, False, INK)
    else:
        for i, creative in enumerate(data["creatives"][:5]):
            _add_textbox(slide, 0.75, 1.0 + i * 0.55, 11.7, 0.32, f"{_date(creative['date'])} | {creative['channel']} | {creative['text'] or '-'}", 12)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, 0.55, 0.35, 5.0, 0.4, "최근 캠페인", 22, True)
    for i, campaign in enumerate(data["campaigns"][:8]):
        _add_textbox(slide, 0.8, 1.05 + i * 0.55, 11.7, 0.32, f"- {campaign['name']} ({campaign['channel']}, {_date(campaign['last_seen'])})", 14)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_pdf(data: dict) -> bytes:
    font = _register_pdf_font()
    styles = getSampleStyleSheet()
    title = ParagraphStyle(
        "AdScopeTitle",
        parent=styles["Title"],
        fontName=font,
        fontSize=22,
        leading=28,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#0f172a"),
    )
    heading = ParagraphStyle("AdScopeHeading", parent=styles["Heading2"], fontName=font, fontSize=14, leading=18)
    body = ParagraphStyle("AdScopeBody", parent=styles["BodyText"], fontName=font, fontSize=10, leading=14)

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, leftMargin=16 * mm, rightMargin=16 * mm, topMargin=16 * mm, bottomMargin=16 * mm)
    story = [
        Paragraph(f"{data['advertiser'].name} 광고 활동 리포트", title),
        Paragraph(f"{_date(data['date_from'])} ~ {_date(data['date_to'])}", body),
        Spacer(1, 8 * mm),
        Paragraph("Executive Summary", heading),
    ]
    for insight in data["summary"]["insights"]:
        story.append(Paragraph(f"- {insight}", body))
    story.append(Spacer(1, 6 * mm))

    kpi = [
        ["광고 소재", "캠페인", "추정 광고비", "주요 채널"],
        [
            f"{data['summary']['total_ads']:,}건",
            f"{data['summary']['campaign_count']:,}건",
            f"{_money(data['summary']['total_spend'])}원",
            data["summary"]["top_channel"],
        ],
    ]
    story.append(Table(kpi, colWidths=[38 * mm] * 4, style=TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0f172a")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, -1), font),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#cbd5e1")),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
        ("TOPPADDING", (0, 0), (-1, -1), 7),
    ])))
    story.append(Spacer(1, 8 * mm))

    story.append(Paragraph("채널별 요약", heading))
    channel_table = [["채널", "소재 수", "최근 관측", "추정 광고비"]]
    for row in data["channels"][:8]:
        channel_table.append([row["channel_label"], f"{row['ad_count']:,}", _date(row["last_seen"]), f"{_money(row['spend'])}원"])
    story.append(Table(channel_table, colWidths=[42 * mm, 28 * mm, 38 * mm, 42 * mm], style=TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#14b8a6")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, -1), font),
        ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#cbd5e1")),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ])))
    story.append(Spacer(1, 8 * mm))

    if data["image_creatives"]:
        story.append(Paragraph("주요 광고 소재 이미지", heading))
        image_rows = []
        for creative in data["image_creatives"][:6]:
            try:
                image = RLImage(creative["image_path"], width=50 * mm, height=36 * mm)
            except Exception:
                image = Paragraph("이미지 로드 실패", body)
            desc = Paragraph(f"{_date(creative['date'])}<br/>{creative['channel']}<br/>{creative['text'] or '-'}", body)
            image_rows.append([image, desc])
        story.append(Table(image_rows, colWidths=[56 * mm, 96 * mm], style=TableStyle([
            ("FONTNAME", (0, 0), (-1, -1), font),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#e2e8f0")),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ("TOPPADDING", (0, 0), (-1, -1), 6),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ])))
        story.append(Spacer(1, 8 * mm))

    story.append(Paragraph("최근 소재", heading))
    for creative in data["creatives"][:8]:
        story.append(Paragraph(f"{_date(creative['date'])} | {creative['channel']} | {creative['text'] or '-'}", body))

    story.append(Spacer(1, 8 * mm))
    story.append(Paragraph("주의사항", heading))
    story.append(Paragraph("본 보고서는 AdScope 수집 데이터와 추정 모델을 기반으로 생성되었습니다. 광고비와 효과 지표는 실제 집행액이 아닌 추정치일 수 있습니다.", body))

    doc.build(story)
    return buf.getvalue()


@router.get("/advertiser/{advertiser_id}/bundle")
async def advertiser_report_bundle(
    advertiser_id: int,
    days: int = Query(30, ge=7, le=365),
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Download a ZIP containing PPTX, PDF, and XLSX advertiser reports."""
    data = await _collect_report_data(db, advertiser_id, days)
    base = f"adscope_{_filename_safe(data['advertiser'].name)}_{_now_kst().strftime('%Y%m%d')}"

    xlsx_bytes = _build_xlsx(data)
    pptx_bytes = _build_pptx(data)
    pdf_bytes = _build_pdf(data)

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(f"{base}.pptx", pptx_bytes)
        zf.writestr(f"{base}.pdf", pdf_bytes)
        zf.writestr(f"{base}.xlsx", xlsx_bytes)
    buf.seek(0)

    filename = f"{base}_bundle.zip"
    return StreamingResponse(
        iter([buf.getvalue()]),
        media_type="application/zip",
        headers=_content_disposition(filename),
    )
